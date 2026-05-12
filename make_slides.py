from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
DARK       = RGBColor(0x1A, 0x1A, 0x1A)
MID        = RGBColor(0x44, 0x44, 0x44)
LIGHT      = RGBColor(0x88, 0x88, 0x88)
RULE       = RGBColor(0xCC, 0xCC, 0xCC)
BLUE       = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
GREEN_BG   = RGBColor(0xE2, 0xEF, 0xDA)
RED_BG     = RGBColor(0xFC, 0xE4, 0xD6)
HEADER_BG  = RGBColor(0x2E, 0x75, 0xB6)
ROW_ALT    = RGBColor(0xF2, 0xF2, 0xF2)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────

def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    return s

def tx(slide, text, x, y, w, h, size=14, bold=False, italic=False,
       color=DARK, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf

def hline(slide, y, x=0.4, w=12.53, color=RULE, thickness=0.5):
    ln = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Pt(thickness))
    ln.fill.solid(); ln.fill.fore_color.rgb = color
    ln.line.fill.background()

def rect(slide, x, y, w, h, fill=WHITE, border=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape

HEADER_ROW = RGBColor(0xE8, 0xE8, 0xE8)

def table_cell(slide, text, x, y, w, h, fill=WHITE, font_size=11,
               bold=False, color=DARK, align=PP_ALIGN.CENTER, border=True,
               is_header=False):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = HEADER_ROW if is_header else fill
    if border:
        shape.line.color.rgb = RULE
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(x + 0.05), Inches(y + 0.02),
                                  Inches(w - 0.1), Inches(h - 0.04))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = True if is_header else bold
    run.font.color.rgb = DARK

def slide_title(slide, title):
    tx(slide, title, 0.4, 0.28, 12.5, 0.55, size=24, bold=True, color=DARK)
    hline(slide, 0.9)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()

hline(s, 1.9, x=0, w=13.33, color=BLUE, thickness=3)

tx(s, "Beyond Answer Generation", 0.6, 2.1, 12.1, 0.9, size=38, bold=True, color=DARK)
tx(s, "Diagnosing Narrative Reasoning in Transformers",
   0.6, 3.0, 12.1, 0.55, size=20, color=MID)

hline(s, 3.75, color=RULE)

tx(s, "W266 — NLP Final Project", 0.6, 3.95, 6, 0.38, size=14, color=LIGHT)
tx(s, "Nityasree Cheera  |  Aurelia Yang", 0.6, 4.33, 6, 0.38, size=14, color=LIGHT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Research Questions
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "Research Questions")

rqs = [
    ("RQ1", "Do transformers truly reason about narrative logic, or do they exploit surface-level patterns such as word overlap and keyword matching?"),
    ("RQ2", "Which type of narrative inconsistency is hardest for models to detect — temporal, causal, state consistency, or sentiment?"),
]

qy = 1.3
for label, q in rqs:
    tx(s, label + ":", 0.4, qy + 0.22, 1.1, 0.45, size=15, bold=True, color=DARK)
    tx(s, q, 1.55, qy + 0.12, 11.2, 0.65, size=15, color=DARK)
    qy += 1.3

tx(s, "Approach: generate controlled negatives for 4 reasoning types, evaluate 6 models across two test sets.",
   0.4, 4.2, 12.5, 0.38, size=13, italic=True, color=LIGHT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Data Description
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "Data Description")

# Datasets table
headers = ["Dataset", "Size", "Purpose"]
col_w   = [3.5, 2.0, 7.0]
col_x   = [0.4, 3.9, 5.9]
row_h   = 0.4

# header row
for i, h in enumerate(headers):
    table_cell(s, h, col_x[i], 1.05, col_w[i], row_h, font_size=12, is_header=True)

rows = [
    ("ROCStories 2017",       "52,665 stories", "Training",  0.38),
    ("Story Cloze Val 2018",  "1,571 stories",  "Testing (Test A & Test B)", 0.38),
    ("Story Cloze Test 2018", "1,571 stories",  "Not used — no labels provided", 0.38),
]
ry = 1.05 + row_h
for ri, (d, sz, purp, rh) in enumerate(rows):
    bg = WHITE if ri % 2 == 0 else ROW_ALT
    table_cell(s, d,    col_x[0], ry, col_w[0], rh, fill=bg, font_size=11, align=PP_ALIGN.LEFT)
    table_cell(s, sz,   col_x[1], ry, col_w[1], rh, fill=bg, font_size=11, align=PP_ALIGN.LEFT)
    table_cell(s, purp, col_x[2], ry, col_w[2], rh, fill=bg, font_size=11, align=PP_ALIGN.LEFT)
    ry += rh

# 4 reasoning types table
tx(s, "4 Controlled Reasoning Types", 0.4, ry + 0.2, 8, 0.35, size=14, bold=True, color=DARK)

neg_headers = ["Type", "Violation"]
neg_col_w   = [2.5, 10.0]
neg_col_x   = [0.4, 2.9]
neg_hdr_y   = ry + 0.62

for i, h in enumerate(neg_headers):
    table_cell(s, h, neg_col_x[i], neg_hdr_y, neg_col_w[i], 0.36, font_size=12, is_header=True)

neg_rows = [
    ("Temporal",  "Events happen in the wrong order, too early, too late, or in the wrong sequence"),
    ("Causal",    "Cause-effect logic is violated; consequences are illogical given the story context"),
    ("State",     "A world-state fact established in the story is contradicted by the ending"),
    ("Sentiment", "The emotional tone of the ending is reversed relative to the story's trajectory"),
]
for ri, (ntype, desc) in enumerate(neg_rows):
    bg = WHITE if ri % 2 == 0 else ROW_ALT
    table_cell(s, ntype, neg_col_x[0], neg_hdr_y + 0.36 + 0.36 * ri, neg_col_w[0], 0.36, fill=bg, font_size=11)
    table_cell(s, desc,  neg_col_x[1], neg_hdr_y + 0.36 + 0.36 * ri, neg_col_w[1], 0.36,
               fill=bg, font_size=11, align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Cloze Task Example
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "The Cloze Task")

tx(s, "Two formats are used: ROC Stories (single correct ending) and Story Cloze Val (2-way forced choice).",
   0.4, 1.0, 12.5, 0.35, size=14, color=MID)

# ROC example
hline(s, 1.45, color=RULE)
tx(s, "ROC Stories", 0.4, 1.52, 3, 0.3, size=12, bold=True, color=DARK)
tx(s, "One correct ending — sentence 5 of the story", 3.5, 1.52, 6, 0.3, size=11, color=LIGHT)

roc_ctx = [
    "S1:  Tom had a very short temper.",
    "S2:  One day a guest made him very angry.",
    "S3:  He punched a hole in the wall of his house.",
    "S4:  Tom's guest became afraid and left quickly.",
]
cy = 1.88
for sent in roc_ctx:
    tx(s, sent, 0.4, cy, 12.5, 0.3, size=13, color=DARK)
    cy += 0.3

tx(s, "S5 (correct ending):  Tom sat on his couch filled with regret about his actions.",
   0.4, cy + 0.05, 12.5, 0.3, size=13, bold=True, color=DARK)

# Cloze Val example
hline(s, 3.65, color=RULE)
tx(s, "Story Cloze Val", 0.4, 3.72, 3, 0.3, size=12, bold=True, color=DARK)
tx(s, "Two human-written endings — pick the correct one", 3.5, 3.72, 6, 0.3, size=11, color=LIGHT)

cloze_ctx = [
    "S1:  Rick grew up in a troubled household.",
    "S2:  He never found good support in family, and turned to gangs.",
    "S3:  It wasn't long before Rick got shot in a robbery.",
    "S4:  The incident caused him to turn a new leaf.",
]
cy2 = 4.08
for sent in cloze_ctx:
    tx(s, sent, 0.4, cy2, 12.5, 0.3, size=13, color=DARK)
    cy2 += 0.3

tx(s, "Option 1:  He is happy now.  (correct)", 0.4, cy2 + 0.05, 12.5, 0.3, size=13, color=DARK)
tx(s, "Option 2:  He joined a gang.  (incorrect)", 0.4, cy2 + 0.35, 12.5, 0.3, size=13, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — EDA
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "Exploratory Data Analysis")

# Left column: ROC
tx(s, "ROCStories", 0.4, 1.05, 5.5, 0.35, size=16, bold=True, color=DARK)

roc_rows = [
    ("Stories",          "52,665"),
    ("Missing values",   "None"),
    ("Avg story length", "~42 words"),
    ("Top bigrams",      "went home, high school, worked hard, got home"),
    ("Topics",           "Family, school, shopping, outdoors"),
]

ry = 1.48
for label, val in roc_rows:
    tx(s, label, 0.4, ry, 2.5, 0.35, size=12, bold=True, color=DARK)
    tx(s, val,   2.95, ry, 3.3, 0.35, size=12, color=MID)
    ry += 0.4

# Right column: Cloze Val
hline(s, 1.05, x=6.6, w=6.33, color=RULE, thickness=0.3)
tx(s, "Story Cloze Val", 6.6, 1.05, 6.1, 0.35, size=16, bold=True, color=DARK)

cloze_rows = [
    ("Stories",        "1,571"),
    ("Missing values", "None"),
    ("Label balance",  "Ending 1 correct: 51.1%   Ending 2: 48.9%"),
    ("Avg ending length", "Option 1: ~7 words   Option 2: ~7 words"),
    ("Shortcut check", "No label or length shortcut exists"),
]

ry2 = 1.48
for label, val in cloze_rows:
    tx(s, label, 6.6,  ry2, 2.6, 0.32, size=12, bold=True, color=DARK)
    tx(s, val,   9.25, ry2, 3.8, 0.32, size=12, color=MID)
    ry2 += 0.38

hline(s, 3.82, color=RULE)
tx(s, "Takeaway: both datasets are balanced with no obvious shortcuts. "
      "Surface cues like word count or label frequency are insufficient to solve the task.",
   0.4, 3.92, 12.5, 0.45, size=13, italic=True, color=MID)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Data Generation
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "Data Generation")

bullets = [
    "ROCStories has no reasoning type labels — so we generate our own incorrect endings using GPT-4o-mini.",
    "We generate 4 wrong endings per ROC story (training) and per Cloze story (evaluation).",
    "Story Cloze is then used in two ways:",
    "   Test A — model picks between two human-written endings (correct vs. incorrect). Tests real-world generalization.",
    "   Test B — model picks between the correct ending and a controlled wrong one. Tests one reasoning type at a time.",
]
by = 1.0
for b in bullets:
    bold = b.startswith("Story Cloze") or b.startswith("   Test")
    tx(s, b, 0.4, by, 12.5, 0.32, size=13, bold=False, color=DARK if not b.startswith("   ") else MID)
    by += 0.33

# Story context
hline(s, by + 0.1, color=RULE)
tx(s, "Story context", 0.4, by + 0.2, 3, 0.3, size=11, bold=True, color=LIGHT)
tx(s, "David noticed he had put on a lot of weight recently. He examined his habits to figure out the reason. "
      "He realized he had been eating too much fast food. He stopped going to burger places and started a vegetarian diet.",
   0.4, by + 0.48, 12.5, 0.45, size=12.5, color=DARK)

tx(s, "Gold ending:  After a few weeks, he started to feel much better.",
   0.4, by + 0.98, 12.5, 0.3, size=12.5, bold=True, color=DARK)

hline(s, by + 1.38, color=RULE)

# Negatives table
tbl_y = by + 1.5
neg_headers2 = ["Type", "Generated Incorrect Ending", "Violation"]
neg_col_w2   = [1.5, 6.8, 4.7]
neg_col_x2   = [0.4, 1.9, 8.7]

for i, h in enumerate(neg_headers2):
    table_cell(s, h, neg_col_x2[i], tbl_y, neg_col_w2[i], 0.36, font_size=12, is_header=True)

gen_rows = [
    ("Temporal",
     "He felt worse immediately after making the change.",
     "Happens before any time has passed — wrong sequence"),
    ("Causal",
     "He found himself gaining even more weight after weeks.",
     "Vegetarian diet causing weight gain — illogical outcome"),
    ("State",
     "After a few weeks, he found himself craving burgers more than ever.",
     "Contradicts that he stopped wanting fast food"),
    ("Sentiment",
     "After a few weeks, he felt sluggish and dissatisfied with his choices.",
     "Emotional tone reversed — story implies a positive arc"),
]

rh2 = 0.46
for ri, (t, ending, violation) in enumerate(gen_rows):
    bg = WHITE if ri % 2 == 0 else ROW_ALT
    table_cell(s, t,         neg_col_x2[0], tbl_y + 0.36 + rh2*ri, neg_col_w2[0], rh2, fill=bg, font_size=11)
    table_cell(s, ending,    neg_col_x2[1], tbl_y + 0.36 + rh2*ri, neg_col_w2[1], rh2, fill=bg, font_size=11, align=PP_ALIGN.LEFT)
    table_cell(s, violation, neg_col_x2[2], tbl_y + 0.36 + rh2*ri, neg_col_w2[2], rh2, fill=bg, font_size=11, align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Data Filtering
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "Data Filtering and Final Splits")

# Filters
tx(s, "Quality filters applied to generated negatives", 0.4, 1.05, 8, 0.35,
   size=15, bold=True, color=DARK)

filters = [
    ("Length",             "Keep endings between 4 and 25 words. Removes empty outputs and multi-sentence responses."),
    ("Sentence count",     "Exactly one sentence per ending. GPT-4o-mini occasionally generates two; these are rejected."),
    ("Uniqueness",         "Generated ending must differ from the gold ending. Prevents near-duplicate negatives."),
    ("Resumability",       "Story IDs are tracked per JSONL bundle. Allows crash-safe generation across 5,000+ stories in batches of 500."),
]

fy = 1.48
for label, desc in filters:
    tx(s, label, 0.4, fy, 2.0, 0.3, size=12.5, bold=True, color=DARK)
    tx(s, desc,  2.45, fy, 10.5, 0.3, size=12.5, color=MID)
    fy += 0.42

hline(s, 3.55, color=RULE)

# Final splits table
tx(s, "Final dataset splits", 0.4, 3.65, 6, 0.35, size=15, bold=True, color=DARK)

split_headers = ["Split", "Size", "Description"]
split_col_w   = [2.8, 1.7, 8.5]
split_col_x   = [0.4, 3.2, 4.9]

for i, h in enumerate(split_headers):
    table_cell(s, h, split_col_x[i], 4.08, split_col_w[i], 0.38, font_size=12, is_header=True)

split_rows = [
    ("Train (NLI pairs)",    "2,512 pairs",   "80% of ROC stories paired with LLM negatives; context + ending -> label"),
    ("Monitor / Val",        "630 pairs",      "20% of ROC; used for early stopping only, never reported"),
    ("Test A (Cloze 2-way)",  "1,571 stories",  "Human-written negatives; generalization benchmark"),
    ("Test B (per-type)",     "~5,493 pairs",   "LLM negatives broken down by reasoning type; diagnostic evaluation"),
]

rh3 = 0.42
for ri, row in enumerate(split_rows):
    bg = WHITE if ri % 2 == 0 else ROW_ALT
    for i, cell in enumerate(row):
        table_cell(s, cell, split_col_x[i], 4.08 + 0.38 + rh3 * ri,
                   split_col_w[i], rh3, fill=bg, font_size=11,
                   align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Models
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "Models")

m_headers = ["Model", "Type", "Description"]
m_col_w   = [3.2, 1.8, 8.0]
m_col_x   = [0.4, 3.6, 5.4]

for i, h in enumerate(m_headers):
    table_cell(s, h, m_col_x[i], 1.05, m_col_w[i], 0.38, font_size=12, is_header=True)

model_rows = [
    ("Lexical Overlap",           "Unsupervised",
     "Counts shared words between context and ending; picks the higher-overlap option"),
    ("Sentence-BERT (MiniLM-L6)", "Unsupervised",
     "Encodes context and ending into 384-dim dense vectors; cosine similarity; no fine-tuning"),
    ("Pairwise DistilBERT",       "Supervised",
     "Fine-tunes DistilBERT on ending pairs; 3 epochs, lr=2e-5, batch 16, early stopping"),
    ("Pairwise RoBERTa-base",     "Supervised",
     "Same pairwise setup as DistilBERT using roberta-base; more pretraining data, dynamic masking"),
]

rh4 = 0.48
for ri, (name, mtype, desc) in enumerate(model_rows):
    bg = WHITE if ri % 2 == 0 else ROW_ALT
    table_cell(s, name,  m_col_x[0], 1.05 + 0.38 + rh4 * ri, m_col_w[0], rh4, fill=bg, font_size=11, align=PP_ALIGN.LEFT)
    table_cell(s, mtype, m_col_x[1], 1.05 + 0.38 + rh4 * ri, m_col_w[1], rh4, fill=bg, font_size=11)
    table_cell(s, desc,  m_col_x[2], 1.05 + 0.38 + rh4 * ri, m_col_w[2], rh4, fill=bg, font_size=11, align=PP_ALIGN.LEFT)

hline(s, 4.4, color=RULE)
tx(s, "Evaluation metric: pairwise accuracy — does the model rank the gold ending above the generated negative?",
   0.4, 4.5, 12.5, 0.35, size=13, italic=True, color=MID)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Final Results
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_title(s, "Final Results")

# Test A — left half (0.4 to 5.8)
tx(s, "Test A — Human-written negatives", 0.4, 1.05, 5.4, 0.35, size=13, bold=True, color=DARK)

vb_col_x = [0.4,  3.6,  4.9]
vb_col_w = [3.15, 1.25, 1.6]
vb_headers = ["Model", "Accuracy", "Significant?"]

for i, h in enumerate(vb_headers):
    table_cell(s, h, vb_col_x[i], 1.45, vb_col_w[i], 0.35, font_size=11, is_header=True)

vb_rows = [
    ("Lexical Overlap",       "51.7%", "No"),
    ("SBERT",                 "59.4%", "Yes"),
    ("DistilBERT (pairwise)", "50.6%", "No"),
    ("RoBERTa (pairwise)",    "57.0%", "Yes"),
]

rh5 = 0.34
for ri, (name, acc, sig) in enumerate(vb_rows):
    bg = WHITE if ri % 2 == 0 else ROW_ALT
    table_cell(s, name, vb_col_x[0], 1.45 + 0.35 + rh5*ri, vb_col_w[0], rh5, fill=bg, font_size=11, align=PP_ALIGN.LEFT)
    table_cell(s, acc,  vb_col_x[1], 1.45 + 0.35 + rh5*ri, vb_col_w[1], rh5, fill=bg, font_size=11)
    table_cell(s, sig,  vb_col_x[2], 1.45 + 0.35 + rh5*ri, vb_col_w[2], rh5, fill=bg, font_size=11)

# Test B — right half (6.8 to 13.2)
tx(s, "Test B — Per-type pairwise accuracy (LLM negatives)", 6.8, 1.05, 6.4, 0.35, size=13, bold=True, color=DARK)

vc_col_x = [6.8,  9.0,  9.85, 10.7, 11.55, 12.45]
vc_col_w = [2.15, 0.82, 0.82, 0.82, 0.87,  0.82]
vc_headers = ["Model", "Temp", "Causal", "State", "Sent", "Overall"]

for i, h in enumerate(vc_headers):
    table_cell(s, h, vc_col_x[i], 1.45, vc_col_w[i], 0.35, font_size=10, is_header=True)

vc_rows = [
    ("Lexical Overlap", "6.6%",  "9.5%",  "11.3%", "13.4%", "10.2%"),
    ("SBERT",           "23.8%", "41.1%", "44.3%", "54.2%", "40.9%"),
    ("DistilBERT",      "99.3%", "99.7%", "98.6%", "98.8%", "99.1%"),
    ("RoBERTa",         "99.7%", "99.9%", "99.5%", "99.7%", "99.7%"),
]

for ri, row in enumerate(vc_rows):
    bg = WHITE if ri % 2 == 0 else ROW_ALT
    for i, cell in enumerate(row):
        al = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER
        table_cell(s, cell, vc_col_x[i], 1.45 + 0.35 + rh5*ri, vc_col_w[i], rh5, fill=bg, font_size=10, align=al)

# Key findings
hline(s, 4.55, color=RULE)
tx(s, "Key findings", 0.4, 4.65, 4, 0.35, size=14, bold=True, color=DARK)

findings = [
    ("Style overfitting",
     "Supervised models score 97-99% on Test B but only 51-59% on Test A. "
     "They learned GPT-4o-mini writing patterns, not story reasoning."),
    ("Temporal is hardest",
     "Unsupervised models score only 24-31% on temporal negatives, "
     "indicating that event-sequence reasoning is beyond surface similarity."),
    ("Sentiment is easiest",
     "Both unsupervised models beat chance only on sentiment negatives (54-55%), "
     "suggesting sentiment cues are more surface-detectable."),
]

fy2 = 5.1
for label, body in findings:
    tx(s, label + ":  " + body, 0.4, fy2, 12.5, 0.38, size=12, color=DARK)
    fy2 += 0.46

# ── Speaker notes ─────────────────────────────────────────────────────────────
notes = [
    # Slide 1 — Title (~15 sec)
    "Hi everyone. Our project is called Beyond Answer Generation — Diagnosing Narrative Reasoning in Transformers. "
    "The core question we're asking is: do transformers actually understand story logic, or are they just matching surface patterns?",

    # Slide 2 — Research Questions (~25 sec)
    "We had two research questions. First, do transformers truly reason about narrative logic, or do they exploit things like "
    "word overlap and keyword matching? Second, which type of narrative inconsistency is hardest for models to detect — "
    "temporal, causal, state, or sentiment violations? To answer these, we generated controlled negatives across all four types "
    "and evaluated four models.",

    # Slide 3 — Data Description (~40 sec)
    "We used two datasets. ROCStories — 52,000 crowd-sourced five-sentence everyday stories — as our training source. "
    "And Story Cloze Val 2018, which we used in two ways: Test A, the original 2-way task with human-written foil endings "
    "to test generalization, and Test B, where we paired the same stories with GPT-4o-mini generated negatives for per-type "
    "diagnostic evaluation. We also had the Cloze test set but couldn't use it — the answer labels are withheld by the organizers.",

    # Slide 4 — The Cloze Task (~35 sec)
    "Quick overview of the two formats. In ROC Stories, sentence 5 is always the correct ending — there's only one. "
    "In the Cloze Val format, you get two human-written candidate endings and must pick the right one. For example: given a "
    "story about Rick turning his life around, option 1 is 'He is happy now' and option 2 is 'He joined a gang.' "
    "The model has to pick correctly.",

    # Slide 5 — EDA (~30 sec)
    "A few key findings from our EDA. ROC Stories average about 42 words per story, with everyday topics like going home, "
    "high school, working hard. The Cloze Val set is well balanced — about 51% correct for option 1 and 49% for option 2. "
    "Ending lengths are nearly identical across options. So there are no easy shortcuts for a model to exploit.",

    # Slide 6 — Data Generation (~40 sec)
    "For each ROC story we used GPT-4o-mini to generate four incorrect endings, one per reasoning type. Take this story about "
    "David switching to a vegetarian diet. The gold ending is 'he started to feel much better.' The temporal negative makes him "
    "feel worse immediately — wrong sequence. The causal negative has him gaining weight from the diet — illogical outcome. "
    "The state negative has him craving burgers — contradicts that he stopped. And the sentiment negative reverses the positive "
    "arc entirely. Each prompt was designed to violate only that one reasoning type.",

    # Slide 7 — Data Filtering (~25 sec)
    "We applied four quality filters: length between 4 and 25 words, exactly one sentence, must differ from the gold ending, "
    "and story IDs were tracked per bundle for crash-safe generation. After filtering we had about 5,500 generated negatives. "
    "The final splits were 2,500 training pairs, a monitor split for early stopping, Test A with 1,571 Cloze stories, "
    "and Test B with the per-type negatives.",

    # Slide 8 — Models (~30 sec)
    "We ran four models. Two unsupervised baselines — Lexical Overlap and Sentence-BERT — which require no training. "
    "And two supervised models — Pairwise DistilBERT and Pairwise RoBERTa — both fine-tuned on ending pairs for 3 epochs. "
    "The evaluation metric throughout is pairwise accuracy: does the model rank the gold ending above the generated negative?",

    # Slide 9 — Final Results (~60 sec)
    "Here are the results. On Test A — human-written negatives — only SBERT at 59% and RoBERTa at 57% generalize above chance "
    "in a statistically significant way. DistilBERT surprisingly collapses to 50%, near random. But on Test B — LLM-generated "
    "negatives — both supervised models hit 99%. That gap is the key finding: supervised models learned GPT-4o-mini's writing "
    "style, not actual story reasoning. The unsupervised models tell a more honest story. SBERT scores only 24% on temporal "
    "negatives but 54% on sentiment — showing that temporal reasoning, which requires understanding event sequences, is far "
    "harder than detecting emotional tone reversals. To summarize: high accuracy on generated negatives does not mean a model "
    "is reasoning. You need human-written evaluation sets and per-type breakdowns to see what's really going on.",
]

for slide, note in zip(prs.slides, notes):
    slide.notes_slide.notes_text_frame.text = note

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/nityasreecheera/Downloads/266.pptx"
prs.save(out)
print(f"Saved: {out}")
